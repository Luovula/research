from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def set_font_style(run, font_name, font_size, bold=False, italic=False, color=None):
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Customize the font style
font_name = "Arial"  # Change to a modern font like Arial, Roboto, or Open Sans
font_size = 12
font_color = RGBColor(40, 40, 40)  # A dark gray color for text

# Slide 1: Introduction
slide1 = presentation.slides.add_slide(presentation.slide_layouts[5])
title1 = slide1.shapes.title
title1.text = "Investigating the Numerical Values of the Cost of Passing in Different Stages of a Go Game"

# Add a separate subtitle shape
txBox = slide1.shapes.add_textbox(Inches(2), Inches(1.5), Inches(12), Inches(1))
subtitle1 = txBox.text_frame
subtitle1.text = "By [Your Name]"

# Apply modern styling to title and subtitle
set_font_style(title1.text_frame.paragraphs[0].runs[0], font_name, 28, bold=True, color=font_color)
set_font_style(subtitle1.paragraphs[0].runs[0], font_name, 18, italic=True, color=font_color)

# Add slides and content with customized font style
slides = []

titles = [
    # Add slide titles as strings here
    "Introduction to Go",
    "Three Stages and Cost of Passing",
    "Monte Carlo Tree Search and Convolutional Neural Networks",
    "Usage of AI in Go: AlphaGo",
    "Previous Research and Data Collection",
    "Hypothesis and Analysis",
    "Conclusion and Future Research",
]

contents = [
    # Add slide content as lists of bullet points here
    ["Why I am interested in this topic", "Motivation for research"],
    ["Stage 1: Opening", "Stage 2: Middle game", "Stage 3: Endgame", "Cost of passing"],
    ["Monte Carlo Tree Search", "Convolutional Neural Networks"],
    ["AlphaGo", "Other AI in Go"],
    ["Review of previous research", "Data collection: 200 games"],
    ["Hypothesis", "Analysis of collected data"],
    ["Conclusion", "Future research"],
]


for i in range(len(titles)):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title = slide.shapes.title
    title.text = titles[i]

    content = slide.placeholders[1].text_frame
    content.clear()  # Remove the default bullet points
    for item in contents[i]:
        paragraph = content.add_paragraph()
        paragraph.space_before = Pt(6)
        run = paragraph.add_run()
        run.text = item
        set_font_style(run, font_name, font_size, color=font_color)

    slides.append(slide)

# Save the presentation
presentation.save("go_game_presentation_modern.pptx")
