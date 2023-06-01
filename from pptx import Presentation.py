from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
presentation = Presentation()

# Slide 1: Introduction
slide1 = presentation.slides.add_slide(presentation.slide_layouts[5])
title1 = slide1.shapes.title
title1.text = "Investigating the Numerical Values of the Cost of Passing in Different Stages of a Go Game"

txBox = slide1.shapes.add_textbox(Inches(2), Inches(1.5), Inches(12), Inches(1))
subtitle1 = txBox.text_frame
subtitle1.text = "By [Your Name]"


# Slide 2: Introduction and motivation
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title2 = slide2.shapes.title
title2.text = "Introduction and Motivation"

content2 = slide2.placeholders[1].text_frame
content2.text = "My interest in this topic..."
content2.add_paragraph().text = "Motivation for research..."

# Slide 3: The Game of Go and its stages
slide3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title3 = slide3.shapes.title
title3.text = "The Game of Go and its Stages"

content3 = slide3.placeholders[1].text_frame
content3.text = "Introduction to the game of Go..."
content3.add_paragraph().text = "Three stages and cost of passing..."

# Slide 4: Monte Carlo Tree Search and Convolutional Neural Networks in Go AI
slide4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title4 = slide4.shapes.title
title4.text = "MCTS and CNNs in Go AI"

content4 = slide4.placeholders[1].text_frame
content4.text = "Definition of Monte Carlo Tree Search..."
content4.add_paragraph().text = "Definition of Convolutional Neural Networks..."
content4.add_paragraph().text = "Usage in Go AI such as AlphaGo..."

# Slide 5: Research Methodology
slide5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title5 = slide5.shapes.title
title5.text = "Research Methodology"

content5 = slide5.placeholders[1].text_frame
content5.text = "Review of previous research..."
content5.add_paragraph().text = "Data collection: cost of passing of 200 games..."

# Slide 6: Hypothesis and Analysis
slide6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title6 = slide6.shapes.title
title6.text = "Hypothesis and Analysis"

content6 = slide6.placeholders[1].text_frame
content6.text = "Hypothesis..."
content6.add_paragraph().text = "Analysis of collected data..."

# Slide 7: Conclusion and Future Research
slide7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title7 = slide7.shapes.title
title7.text = "Conclusion and Future Research"

content7 = slide7.placeholders[1].text_frame
content7.text = "Conclusion..."
content7.add_paragraph().text = "Future research directions..."

# Save the presentation
presentation.save("go_game_presentation.pptx")
